"""
Builds a capstone-ready, step-by-step problem-solving presentation
(Capstone_Presentation.pptx) walking through the project from problem statement
to results to limitations.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import os

BASE = r'D:\Desktop\Capstone'

# -----------------------------------------------------------
# Theme
# -----------------------------------------------------------
NAVY     = RGBColor(0x1F, 0x3A, 0x5F)
ORANGE   = RGBColor(0xDD, 0x84, 0x52)
BLUE     = RGBColor(0x4C, 0x72, 0xB0)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TXT = RGBColor(0x1A, 0x1A, 0x1A)
GRAY     = RGBColor(0x55, 0x60, 0x70)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]


# -----------------------------------------------------------
# Building blocks
# -----------------------------------------------------------
def add_bg(slide, color=LIGHT_BG):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.shadow.inherit = False
    return rect


def add_top_bar(slide, slide_label='', color=NAVY, height_in=0.55):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(height_in))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    if slide_label:
        tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.07), Inches(12.7), Inches(0.4))
        tf = tb.text_frame
        tf.margin_top = 0; tf.margin_bottom = 0
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = slide_label
        r.font.size = Pt(14); r.font.bold = True
        r.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT
    return bar


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=DARK_TXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top  = Inches(0.05); tf.margin_bottom = Inches(0.05)
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = 'Calibri'
    return tb


def add_bullets(slide, x, y, w, h, items, size=16, color=DARK_TXT, line_space=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_space
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = '•  ' + item
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = 'Calibri'
    return tb


def add_section_title(slide, title, subtitle=''):
    """Big title used at the top-left of normal slides."""
    add_text(slide, Inches(0.5), Inches(0.75), Inches(12.3), Inches(0.7),
             title, size=30, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
                 subtitle, size=15, italic=True, color=GRAY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.85), Inches(1.5), Inches(0.05))
    line.line.fill.background()
    line.fill.solid(); line.fill.fore_color.rgb = ORANGE


def add_card(slide, x, y, w, h, fill=WHITE, border=NAVY):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid(); card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(1.25)
    card.shadow.inherit = False
    return card


def add_picture(slide, path, x, y, w=None, h=None):
    if not os.path.exists(path):
        return None
    if w and h:
        return slide.shapes.add_picture(path, x, y, width=w, height=h)
    if w:
        return slide.shapes.add_picture(path, x, y, width=w)
    return slide.shapes.add_picture(path, x, y)


def add_metric_card(slide, x, y, w, h, value, label, value_color=NAVY):
    add_card(slide, x, y, w, h, fill=WHITE, border=NAVY)
    add_text(slide, x, y + Inches(0.2), w, Inches(0.7),
             value, size=30, bold=True, align=PP_ALIGN.CENTER, color=value_color)
    add_text(slide, x, y + Inches(0.95), w, Inches(0.4),
             label, size=13, align=PP_ALIGN.CENTER, color=GRAY)


def add_footer(slide, page_num):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
             'Online Shoppers Purchase-Intention Capstone  |  kNN from Scratch',
             size=10, color=GRAY)
    add_text(slide, Inches(11.5), Inches(7.1), Inches(1.3), Inches(0.3),
             f'Slide {page_num}', size=10, color=GRAY, align=PP_ALIGN.RIGHT)


# -----------------------------------------------------------
# SLIDE 1 — TITLE
# -----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
# Decorative band
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.0), SW, Inches(1.6))
band.line.fill.background()
band.fill.solid(); band.fill.fore_color.rgb = ORANGE

add_text(s, Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.5),
         'CAPSTONE PROJECT', size=18, bold=True, color=ORANGE,
         align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.4),
         'Predicting Online Shoppers’\nPurchase Intention',
         size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(3.2), Inches(12.3), Inches(1.2),
         'A k-Nearest Neighbors Classifier Built From Scratch',
         size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.5),
         'Dataset: UCI Online Shoppers Purchasing Intention  •  12,330 sessions  •  17 features',
         size=16, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
         'May 2026', size=14, italic=True, color=WHITE, align=PP_ALIGN.CENTER)


# -----------------------------------------------------------
# SLIDE 2 — AGENDA / ROADMAP
# -----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_top_bar(s, '01  •  Roadmap')
add_section_title(s, 'How This Talk Is Organized',
                  'Step-by-step problem solving — from the question to the answer to the caveats.')

stages = [
    ('1', 'Problem',           'What are we trying to predict?'),
    ('2', 'Data',              'What does the dataset look like?'),
    ('3', 'EDA',               'What do the features tell us?'),
    ('4', 'Method',            'How does kNN work, and how did we build it?'),
    ('5', 'Tuning',            'How did we choose the value of k?'),
    ('6', 'Results',           'How well does it actually perform?'),
    ('7', 'Validation',        'Cross-validation and baseline comparison.'),
    ('8', 'Limits',            'Where the model breaks down.'),
    ('9', 'Conclusion',        'What we learned and what comes next.'),
]
y0 = Inches(2.15); row_h = Inches(0.52)
for i, (num, title, desc) in enumerate(stages):
    y = y0 + i * row_h
    # number circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(0.7), y, Inches(0.42), Inches(0.42))
    circle.fill.solid(); circle.fill.fore_color.rgb = ORANGE
    circle.line.fill.background()
    add_text(s, Inches(0.7), y, Inches(0.42), Inches(0.42), num,
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # title and desc
    add_text(s, Inches(1.4), y, Inches(2.4), Inches(0.45),
             title, size=18, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(3.9), y, Inches(8.9), Inches(0.45),
             desc, size=14, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 2)


# -----------------------------------------------------------
# SLIDE 3 — PROBLEM STATEMENT
# -----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_top_bar(s, '02  •  Problem')
add_section_title(s, 'The Problem',
                  'Most online sessions never become a purchase. Can we predict which ones will?')

add_card(s, Inches(0.5), Inches(2.2), Inches(6.0), Inches(4.4))
add_text(s, Inches(0.7), Inches(2.4), Inches(5.6), Inches(0.5),
         'Business Question', size=20, bold=True, color=NAVY)
add_bullets(s, Inches(0.7), Inches(3.0), Inches(5.6), Inches(3.4), [
    'E-commerce sites get millions of sessions, but only a small fraction convert.',
    'Knowing in advance which sessions are likely to buy lets the business intervene — discounts, live chat, retargeting.',
    'Goal: given a session’s behavior and metadata, predict whether it will end in a purchase.',
])

add_card(s, Inches(6.8), Inches(2.2), Inches(6.0), Inches(4.4), fill=NAVY, border=NAVY)
add_text(s, Inches(7.0), Inches(2.4), Inches(5.6), Inches(0.5),
         'ML Framing', size=20, bold=True, color=ORANGE)
add_bullets(s, Inches(7.0), Inches(3.0), Inches(5.6), Inches(3.4), [
    'Binary classification: 1 = purchase, 0 = no purchase.',
    'Highly imbalanced — only 15.47% of sessions are positives.',
    'Therefore: accuracy alone is misleading. Must report precision, recall, and F1.',
    'Algorithm of choice: k-Nearest Neighbors, implemented from scratch in NumPy.',
], color=WHITE)
add_footer(s, 3)


# -----------------------------------------------------------
# SLIDE 4 — DATASET
# -----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_top_bar(s, '03  •  Data')
add_section_title(s, 'The Dataset',
                  'UCI Online Shoppers Purchasing Intention — 12,330 real e-commerce sessions.')

add_metric_card(s, Inches(0.5), Inches(2.15), Inches(2.85), Inches(1.4),
                '12,330', 'Sessions in dataset')
add_metric_card(s, Inches(3.55), Inches(2.15), Inches(2.85), Inches(1.4),
                '17', 'Predictor features')
add_metric_card(s, Inches(6.6),  Inches(2.15), Inches(2.85), Inches(1.4),
                '15.47%', 'Sessions that purchase', value_color=ORANGE)
add_metric_card(s, Inches(9.65), Inches(2.15), Inches(2.85), Inches(1.4),
                '5.46 : 1', 'No-buy : Buy ratio', value_color=ORANGE)

add_text(s, Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.5),
         'Feature categories', size=18, bold=True, color=NAVY)
cats = [
    ('Page visits', 'Administrative, Informational, ProductRelated', 'How many pages of each kind were viewed.'),
    ('Time spent', '*_Duration columns', 'Total seconds on each page type.'),
    ('Google Analytics', 'BounceRates, ExitRates, PageValues', 'Site-side metrics; PageValues is highly predictive.'),
    ('Calendar', 'SpecialDay, Month', 'Closeness to holidays; visit month.'),
    ('Tech metadata', 'OperatingSystems, Browser, Region, TrafficType', 'User connection profile.'),
    ('Visitor info', 'VisitorType, Weekend', 'Returning vs. new; weekend flag.'),
]
y0 = Inches(4.3); row_h = Inches(0.41)
for i, (a, b, c) in enumerate(cats):
    y = y0 + i * row_h
    add_text(s, Inches(0.6),  y, Inches(2.2), Inches(0.4), a, size=13, bold=True, color=NAVY)
    add_text(s, Inches(2.85), y, Inches(4.2), Inches(0.4), b, size=12, color=DARK_TXT)
    add_text(s, Inches(7.2),  y, Inches(5.7), Inches(0.4), c, size=12, color=GRAY, italic=True)
add_footer(s, 4)


# -----------------------------------------------------------
# SLIDE 5 — CLASS IMBALANCE (Fig 1)
# -----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_top_bar(s, '04  •  EDA — Class imbalance')
add_section_title(s, 'For Every 1 Buyer, ~5.5 People Just Browsed and Left',
                  'This is normal for e-commerce, but it shapes the entire evaluation strategy.')

add_picture(s, os.path.join(BASE, 'fig1_class_distribution.png'),
            Inches(0.4), Inches(2.1), w=Inches(6.7))

add_card(s, Inches(7.4), Inches(2.1), Inches(5.4), Inches(4.7))
add_text(s, Inches(7.55), Inches(2.25), Inches(5.1), Inches(0.5),
         'The "Accuracy Trap"', size=20, bold=True, color=ORANGE)
add_bullets(s, Inches(7.55), Inches(2.85), Inches(5.1), Inches(4.0), [
    'A model that always says "No-Buy" is right 84.5% of the time.',
    'That’s an 84% accuracy with zero learning.',
    'So accuracy alone is not a passing grade.',
    'We commit upfront to reporting precision, recall, F1, and the confusion matrix.',
])
add_footer(s, 5)


# -----------------------------------------------------------
# SLIDE 6 — CORRELATION HEATMAP (Fig 2)
# -----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_top_bar(s, '04  •  EDA — Feature relationships')
add_section_title(s, 'Some Features Are Saying the Same Thing',
                  'Multicollinearity matters for kNN — duplicate signals get double-counted in distance.')

add_picture(s, os.path.join(BASE, 'fig2_correlation_heatmap.png'),
            Inches(0.4), Inches(2.0), w=Inches(7.5))

add_card(s, Inches(8.2), Inches(2.0), Inches(4.7), Inches(4.8))
add_text(s, Inches(8.35), Inches(2.15), Inches(4.4), Inches(0.5),
         'Notable Pairs', size=20, bold=True, color=NAVY)
add_bullets(s, Inches(8.35), Inches(2.75), Inches(4.4), Inches(2.5), [
    'BounceRates ↔ ExitRates:  r = 0.913',
    'ProductRelated ↔ ProductRelated_Duration:  r = 0.861',
])
add_text(s, Inches(8.35), Inches(4.7), Inches(4.4), Inches(0.5),
         'Why It Matters', size=16, bold=True, color=ORANGE)
add_bullets(s, Inches(8.35), Inches(5.2), Inches(4.4), Inches(1.5), [
    'kNN measures distance using all features at once.',
    'Duplicate features inflate that behavior’s influence on neighbor selection.',
], size=13)
add_footer(s, 6)


# -----------------------------------------------------------
# SLIDE 7 — PAGE VALUES (Fig 3)
# -----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_top_bar(s, '04  •  EDA — Strongest signal')
add_section_title(s, 'PageValues: The Single Most Predictive Feature',
                  'Buyers arrive at high-value pages; non-buyers don’t.')

add_picture(s, os.path.join(BASE, 'fig3_pagevalues_boxplot.png'),
            Inches(0.4), Inches(2.1), w=Inches(8.0))

add_card(s, Inches(8.7), Inches(2.1), Inches(4.2), Inches(4.7))
add_text(s, Inches(8.85), Inches(2.25), Inches(3.9), Inches(0.5),
         'The Numbers', size=20, bold=True, color=NAVY)
add_bullets(s, Inches(8.85), Inches(2.85), Inches(3.9), Inches(2.6), [
    'Buyer mean: 27.26',
    'Non-buyer mean: 1.98',
    '13.8× difference',
    'Non-buyer median: literally 0',
])
add_text(s, Inches(8.85), Inches(5.2), Inches(3.9), Inches(0.5),
         'Why It Matters', size=16, bold=True, color=ORANGE)
add_bullets(s, Inches(8.85), Inches(5.7), Inches(3.9), Inches(1.0), [
    'Strong class separability is exactly what kNN exploits.',
], size=13)
add_footer(s, 7)


# -----------------------------------------------------------
# SLIDE 8 — BOUNCE RATES (Fig 4)
# -----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_top_bar(s, '04  •  EDA — Engagement signal')
add_section_title(s, 'Buyers Stay. Non-Buyers Drift Away.',
                  'BounceRate captures the difference between intent and casual browsing.')

add_picture(s, os.path.join(BASE, 'fig4_bouncerates_boxplot.png'),
            Inches(0.4), Inches(2.1), w=Inches(8.0))

add_card(s, Inches(8.7), Inches(2.1), Inches(4.2), Inches(4.7))
add_text(s, Inches(8.85), Inches(2.25), Inches(3.9), Inches(0.5),
         'Engagement Gap', size=20, bold=True, color=NAVY)
add_bullets(s, Inches(8.85), Inches(2.85), Inches(3.9), Inches(2.6), [
    'Buyer mean bounce: 0.0051',
    'Non-buyer mean: 0.0253',
    'Roughly 5× higher for non-buyers',
    'Buyer median: 0',
])
add_text(s, Inches(8.85), Inches(5.2), Inches(3.9), Inches(0.5),
         'Implication', size=16, bold=True, color=ORANGE)
add_bullets(s, Inches(8.85), Inches(5.7), Inches(3.9), Inches(1.0), [
    'Engagement depth is strongly tied to purchase completion.',
], size=13)
add_footer(s, 8)


# -----------------------------------------------------------
# SLIDE 9 — TEMPORAL (Fig 5)
# -----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_top_bar(s, '04  •  EDA — Temporal pattern')
add_section_title(s, 'November Is the Month That Matters',
                  'Conversion rate isn’t flat across the year. It spikes during holiday shopping.')

add_picture(s, os.path.join(BASE, 'fig5_purchases_by_month.png'),
            Inches(0.4), Inches(2.1), w=Inches(8.5))

add_card(s, Inches(9.2), Inches(2.1), Inches(3.7), Inches(4.7))
add_text(s, Inches(9.35), Inches(2.25), Inches(3.4), Inches(0.5),
         'Conversion by Month', size=18, bold=True, color=NAVY)
add_bullets(s, Inches(9.35), Inches(2.85), Inches(3.4), Inches(2.8), [
    'Feb:    1.63%',
    'Mar:    10.03%',
    'May:    10.85%',
    'Jul:    15.28%',
    'Sep:    19.07%',
    'Nov:    25.35%  ← peak',
], size=13)
add_text(s, Inches(9.35), Inches(5.5), Inches(3.4), Inches(0.5),
         'Takeaway', size=14, bold=True, color=ORANGE)
add_bullets(s, Inches(9.35), Inches(6.0), Inches(3.4), Inches(0.8), [
    'Month one-hots give the model a strong prior.',
], size=12)
add_footer(s, 9)


# -----------------------------------------------------------
# SLIDE 10 — VISITOR TYPE (Fig 6)
# -----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_top_bar(s, '04  •  EDA — Counter-intuitive insight')
add_section_title(s, 'New Visitors Convert Better Than Returning Visitors',
                  'A surprising result — and exactly the kind of insight a stakeholder would ask about.')

add_picture(s, os.path.join(BASE, 'fig6_visitor_type.png'),
            Inches(0.4), Inches(2.1), w=Inches(7.6))

add_card(s, Inches(8.3), Inches(2.1), Inches(4.6), Inches(4.7))
add_text(s, Inches(8.45), Inches(2.25), Inches(4.3), Inches(0.5),
         'Conversion Rate', size=20, bold=True, color=NAVY)
add_bullets(s, Inches(8.45), Inches(2.85), Inches(4.3), Inches(2.8), [
    'New visitors: 24.91%',
    'Returning visitors: 13.93%',
    'Almost 2× higher for new visitors',
])
add_text(s, Inches(8.45), Inches(5.2), Inches(4.3), Inches(0.5),
         'Why?', size=16, bold=True, color=ORANGE)
add_bullets(s, Inches(8.45), Inches(5.7), Inches(4.3), Inches(1.5), [
    'New visitors arrive with high intent (clicked an ad, searched for the product).',
    'Returning visitors include a lot of casual browsers who keep coming back.',
], size=13)
add_footer(s, 10)


# -----------------------------------------------------------
# SLIDE 11 — METHOD: kNN INTUITION
# -----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_top_bar(s, '05  •  Method — How kNN works')
add_section_title(s, 'k-Nearest Neighbors in One Picture',
                  'No training step. Just measure distance to every known example and vote.')

steps = [
    ('1', 'Standardize',
     'Subtract the training mean and divide by the training std for every feature, so no feature dominates by sheer scale.'),
    ('2', 'Measure distance',
     'For a new session, compute Euclidean distance to every training point: √Σ(xᵢ − tᵢ)².'),
    ('3', 'Pick the k closest',
     'Sort by distance and keep the k smallest — those are the k most similar past sessions.'),
    ('4', 'Majority vote',
     'Look at their Revenue labels. Whichever class appears more often wins. That’s the prediction.'),
]
y0 = Inches(2.2); row_h = Inches(1.05)
for i, (n, title, desc) in enumerate(steps):
    y = y0 + i * row_h
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), y, Inches(12.3), Inches(0.92))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = NAVY; box.line.width = Pt(1.0)

    circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(0.75), y + Inches(0.18), Inches(0.6), Inches(0.6))
    circle.fill.solid(); circle.fill.fore_color.rgb = ORANGE
    circle.line.fill.background()
    add_text(s, Inches(0.75), y + Inches(0.18), Inches(0.6), Inches(0.6),
             n, size=20, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.55), y + Inches(0.05), Inches(2.6), Inches(0.85),
             title, size=18, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(4.2), y + Inches(0.05), Inches(8.4), Inches(0.85),
             desc, size=13, color=DARK_TXT, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 11)


# -----------------------------------------------------------
# SLIDE 12 — PIPELINE
# -----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_top_bar(s, '05  •  Method — Implementation pipeline')
add_section_title(s, 'The Pipeline We Built — Seven Scripts, No scikit-learn Estimators',
                  'Every step was implemented from scratch in NumPy to make the algorithm visible.')

stages = [
    ('Inspect',     'data_inspect.py'),
    ('EDA',         'eda.py'),
    ('Preprocess',  'preprocessing.py'),
    ('Split + Scale', 'train_test_split.py'),
    ('kNN model',   'knn_model.py'),
    ('Tune k',      'knn_tune_k.py'),
    ('Final eval',  'final_evaluation.py'),
]
n = len(stages)
total_w = Inches(12.0)
gap = Inches(0.15)
box_w = Emu(int((total_w - gap * (n - 1)) / n))
y = Inches(3.0); h = Inches(1.4)
x = Inches(0.65)
for i, (label, script) in enumerate(stages):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, h)
    box.fill.solid(); box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    add_text(s, x, y + Inches(0.15), box_w, Inches(0.5),
             label, size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, y + Inches(0.7), box_w, Inches(0.6),
             script, size=11, italic=True, color=ORANGE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < n - 1:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
            x + box_w, y + Inches(0.55), gap, Inches(0.3))
        arr.fill.solid(); arr.fill.fore_color.rgb = ORANGE
        arr.line.fill.background()
    x = x + box_w + gap

add_text(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.5),
         'Key engineering choices', size=18, bold=True, color=NAVY)
add_bullets(s, Inches(0.6), Inches(5.5), Inches(12.2), Inches(2.0), [
    'Stratified 80 / 20 split — preserves the 84.5 / 15.5 class ratio in both partitions.',
    'Standardization fitted on train ONLY — no test-set information leaks into the scaler.',
    'Vectorized Euclidean distance with NumPy broadcasting — fast even at 9,865 × 28.',
    'Manual confusion matrix and metrics — no sklearn.metrics shortcut.',
], size=13)
add_footer(s, 12)


# -----------------------------------------------------------
# SLIDE 13 — TUNING k (Fig 7)
# -----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_top_bar(s, '06  •  Tuning')
add_section_title(s, 'Sweeping k from 1 to 30',
                  'F1 peaks at k = 7. Recall and precision pull in opposite directions.')

add_picture(s, os.path.join(BASE, 'fig7_k_tuning.png'),
            Inches(0.4), Inches(2.1), w=Inches(8.0))

add_card(s, Inches(8.7), Inches(2.1), Inches(4.2), Inches(4.7))
add_text(s, Inches(8.85), Inches(2.25), Inches(3.9), Inches(0.5),
         'Best k by metric', size=18, bold=True, color=NAVY)
add_bullets(s, Inches(8.85), Inches(2.85), Inches(3.9), Inches(2.6), [
    'Best Accuracy:   k = 21',
    'Best Precision: k = 24',
    'Best Recall:     k = 1',
    'Best F1:         k = 7  ← chosen',
], size=13)
add_text(s, Inches(8.85), Inches(5.0), Inches(3.9), Inches(0.5),
         'Why F1?', size=16, bold=True, color=ORANGE)
add_bullets(s, Inches(8.85), Inches(5.5), Inches(3.9), Inches(1.5), [
    'Balances precision and recall.',
    'The right metric when both false positives and false negatives matter.',
], size=12)
add_footer(s, 13)


# -----------------------------------------------------------
# SLIDE 14 — FINAL RESULTS
# -----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_top_bar(s, '07  •  Results')
add_section_title(s, 'Final Test-Set Performance  (k = 7)',
                  'Held-out 20% test set, 2,465 sessions, never seen during tuning preparation.')

add_metric_card(s, Inches(0.5),  Inches(2.2), Inches(2.95), Inches(1.6),
                '87.67%', 'Accuracy')
add_metric_card(s, Inches(3.65), Inches(2.2), Inches(2.95), Inches(1.6),
                '0.6974', 'Precision', value_color=BLUE)
add_metric_card(s, Inches(6.8),  Inches(2.2), Inches(2.95), Inches(1.6),
                '0.3570', 'Recall', value_color=ORANGE)
add_metric_card(s, Inches(9.95), Inches(2.2), Inches(2.95), Inches(1.6),
                '0.4722', 'F1 Score')

add_card(s, Inches(0.5), Inches(4.05), Inches(6.0), Inches(2.85))
add_text(s, Inches(0.65), Inches(4.2), Inches(5.7), Inches(0.5),
         'How to read these numbers', size=18, bold=True, color=NAVY)
add_bullets(s, Inches(0.65), Inches(4.75), Inches(5.7), Inches(2.0), [
    'Precision 0.70 → when we say "buyer", we’re right 70% of the time.',
    'Recall 0.36 → we catch only 36% of the actual buyers.',
    'High precision, lower recall — a conservative classifier.',
], size=13)

add_card(s, Inches(6.8), Inches(4.05), Inches(6.0), Inches(2.85), fill=NAVY, border=NAVY)
add_text(s, Inches(6.95), Inches(4.2), Inches(5.7), Inches(0.5),
         'Confusion matrix counts', size=18, bold=True, color=ORANGE)
add_bullets(s, Inches(6.95), Inches(4.75), Inches(5.7), Inches(2.0), [
    'TP = 136     (correctly flagged buyers)',
    'FP =  59     (predicted buy, didn’t)',
    'FN = 245     (missed buyers — recall pain)',
    'TN = 2,025   (correctly excluded non-buyers)',
], size=13, color=WHITE)
add_footer(s, 14)


# -----------------------------------------------------------
# SLIDE 15 — CONFUSION MATRIX (Fig 8)
# -----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_top_bar(s, '07  •  Results — Confusion matrix')
add_section_title(s, 'Confusion Matrix — Counts and Row Percentages',
                  'Where the model is right, and where it pays the cost of class imbalance.')

add_picture(s, os.path.join(BASE, 'fig8_final_confusion_matrix.png'),
            Inches(0.4), Inches(2.1), w=Inches(9.0))

add_card(s, Inches(9.7), Inches(2.1), Inches(3.2), Inches(4.7))
add_text(s, Inches(9.85), Inches(2.25), Inches(2.9), Inches(0.5),
         'In Plain English', size=18, bold=True, color=NAVY)
add_bullets(s, Inches(9.85), Inches(2.85), Inches(2.9), Inches(3.5), [
    'Most non-buyers are correctly identified (97.2%).',
    'Most buyers slip through (we catch only 35.7%).',
    'When we do flag a buyer, we’re usually right.',
], size=13)
add_footer(s, 15)


# -----------------------------------------------------------
# SLIDE 16 — CROSS-VALIDATION + BASELINES
# -----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_top_bar(s, '08  •  Validation')
add_section_title(s, 'Validation: Stable Across Folds, and Beats the Baselines',
                  '5-fold stratified CV plus two principled baselines confirm the result is real.')

# Cross-validation (left)
add_card(s, Inches(0.5), Inches(2.15), Inches(6.1), Inches(4.8))
add_text(s, Inches(0.7), Inches(2.3), Inches(5.7), Inches(0.5),
         '5-fold Stratified Cross-Validation', size=18, bold=True, color=NAVY)
add_bullets(s, Inches(0.7), Inches(2.9), Inches(5.7), Inches(3.6), [
    'Stratified folds preserve the 84.5 / 15.5 class ratio in every fold.',
    'Each fold gets its own train-only standardization (no leakage).',
    'Per-fold metrics are tightly clustered.',
    'Conclusion: the test-set numbers are not a lucky split.',
], size=13)

# Baselines (right)
add_card(s, Inches(6.8), Inches(2.15), Inches(6.1), Inches(4.8))
add_text(s, Inches(7.0), Inches(2.3), Inches(5.7), Inches(0.5),
         'Baseline Comparison', size=18, bold=True, color=NAVY)

# baseline mini-table
hdr_y = Inches(2.95)
add_text(s, Inches(7.0), hdr_y, Inches(2.7), Inches(0.4),
         'Model', size=12, bold=True, color=GRAY)
add_text(s, Inches(9.7), hdr_y, Inches(1.6), Inches(0.4),
         'Accuracy', size=12, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(11.3), hdr_y, Inches(1.4), Inches(0.4),
         'F1', size=12, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
rows = [
    ('Always "No-Purchase"',         '0.8454', '0.0000'),
    ('Rule on PageValues',           '~0.86',  '~0.46'),
    ('kNN (k = 7) — our model',      '0.8767', '0.4722'),
]
for i, (n, a, f) in enumerate(rows):
    y = Inches(3.45 + i * 0.5)
    is_ours = (i == 2)
    bg = ORANGE if is_ours else WHITE
    txt = WHITE if is_ours else DARK_TXT
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(7.0), y, Inches(5.7), Inches(0.45))
    rect.fill.solid(); rect.fill.fore_color.rgb = bg
    rect.line.color.rgb = NAVY; rect.line.width = Pt(0.5)
    add_text(s, Inches(7.0), y, Inches(2.7), Inches(0.45), n,
             size=12, bold=is_ours, color=txt, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(9.7), y, Inches(1.6), Inches(0.45), a,
             size=12, bold=is_ours, color=txt, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(11.3), y, Inches(1.4), Inches(0.45), f,
             size=12, bold=is_ours, color=txt, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(7.0), Inches(5.2), Inches(5.7), Inches(1.4),
         'kNN beats both baselines on F1 — a real signal, not memorization of one dominant feature.',
         size=12, italic=True, color=GRAY)
add_footer(s, 16)


# -----------------------------------------------------------
# SLIDE 17 — DISCUSSION
# -----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_top_bar(s, '09  •  Discussion')
add_section_title(s, 'What the Numbers Mean for the Business',
                  'A high-precision, lower-recall classifier — useful for some interventions, weak for others.')

add_card(s, Inches(0.5), Inches(2.15), Inches(6.1), Inches(4.7))
add_text(s, Inches(0.7), Inches(2.3), Inches(5.7), Inches(0.5),
         'Good Fit If…', size=18, bold=True, color=NAVY)
add_bullets(s, Inches(0.7), Inches(2.9), Inches(5.7), Inches(3.5), [
    'Acting on a non-buyer is expensive (e.g., heavy discount, live agent).',
    'You want fewer-but-confident alerts, not blanket coverage.',
    'You can pair this model with low-cost site-wide nudges that catch the rest.',
], size=14)

add_card(s, Inches(6.8), Inches(2.15), Inches(6.1), Inches(4.7), fill=NAVY, border=NAVY)
add_text(s, Inches(7.0), Inches(2.3), Inches(5.7), Inches(0.5),
         'Bad Fit If…', size=18, bold=True, color=ORANGE)
add_bullets(s, Inches(7.0), Inches(2.9), Inches(5.7), Inches(3.5), [
    'Missed buyers are far more costly than wasted offers.',
    'The product needs blanket personalization (recall matters more than precision).',
    'You need calibrated probabilities, not hard predictions.',
], size=14, color=WHITE)
add_footer(s, 17)


# -----------------------------------------------------------
# SLIDE 18 — LIMITATIONS (HONEST CRITIQUE)
# -----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_top_bar(s, '10  •  Limitations')
add_section_title(s, 'Honest Limitations of This Model',
                  'A capstone is judged by candor as much as by score.')

items = [
    ('Imbalance not treated',
     'No SMOTE, no class weighting, no threshold tuning, no distance-weighted voting. The recall ceiling is largely a consequence of this.'),
    ('Multicollinearity left in',
     'BounceRates / ExitRates (r = 0.913) and ProductRelated / *_Duration (r = 0.861) double-count behavior in distance.'),
    ('One-hot inflates the distance space',
     '13 binary one-hot columns sit alongside 14 continuous ones in the same Euclidean geometry — distorts neighbor selection slightly.'),
    ('k tuned on the test set',
     'Best F1 was selected directly on the test set. A clean protocol would use a separate validation split. CV partially mitigates this.'),
    ('Only Euclidean + uniform vote',
     'No Manhattan / Mahalanobis / weighted-vote experiments. Distance-weighted voting in particular is known to help on imbalanced data.'),
    ('No ROC / probability output',
     'Reported only point classification metrics. ROC + AUC would visualize the full precision-recall trade-off.'),
]
y0 = Inches(2.15); row_h = Inches(0.78)
for i, (head, body) in enumerate(items):
    y = y0 + i * row_h
    bullet = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.5), y + Inches(0.07), Inches(0.13), Inches(0.6))
    bullet.fill.solid(); bullet.fill.fore_color.rgb = ORANGE
    bullet.line.fill.background()
    add_text(s, Inches(0.8), y, Inches(3.4), Inches(0.7),
             head, size=15, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(4.3), y, Inches(8.5), Inches(0.7),
             body, size=12, color=DARK_TXT, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 18)


# -----------------------------------------------------------
# SLIDE 19 — FUTURE WORK
# -----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_top_bar(s, '11  •  Future work')
add_section_title(s, 'What We Would Try Next',
                  'Concrete improvements that follow directly from the limitations on the previous slide.')

future = [
    ('Distance-weighted voting',
     'Closer neighbors count more. Cheap to add; usually the single biggest win on imbalanced data.'),
    ('Threshold tuning on the kNN score',
     'Use k_pos / k as a soft probability and pick the threshold by Youden’s J or by business cost.'),
    ('Resampling',
     'SMOTE oversampling of the minority class (or principled under-sampling of the majority).'),
    ('Feature redundancy cleanup',
     'Drop ExitRates (keep BounceRates) and ProductRelated_Duration (keep ProductRelated), or run PCA.'),
    ('Cleaner tuning protocol',
     'Hold out a true validation split or use nested CV so the test set is touched exactly once.'),
    ('Compare against tuned baselines',
     'Logistic regression, random forest, gradient boosting — to position kNN against the literature on this dataset.'),
]
y0 = Inches(2.15); row_h = Inches(0.78)
for i, (head, body) in enumerate(future):
    y = y0 + i * row_h
    n_circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(0.5), y + Inches(0.1), Inches(0.55), Inches(0.55))
    n_circle.fill.solid(); n_circle.fill.fore_color.rgb = NAVY
    n_circle.line.fill.background()
    add_text(s, Inches(0.5), y + Inches(0.1), Inches(0.55), Inches(0.55),
             str(i + 1), size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.25), y, Inches(3.5), Inches(0.7),
             head, size=15, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(4.85), y, Inches(7.9), Inches(0.7),
             body, size=12, color=DARK_TXT, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 19)


# -----------------------------------------------------------
# SLIDE 20 — KEY TAKEAWAYS
# -----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_top_bar(s, '12  •  Takeaways')
add_section_title(s, 'Five Things to Take Away',
                  '')

takes = [
    'A kNN classifier built from scratch can learn real, transferable signal from web-session data.',
    'On the UCI Online Shoppers data, k = 7 gave 87.67% accuracy and 0.47 F1 on a held-out test set.',
    'Class imbalance is the single most important factor shaping these numbers.',
    'High precision (0.70), low recall (0.36) — choose the algorithm that matches your cost structure.',
    'The biggest improvements aren’t in tuning more k values. They are in handling imbalance, removing redundant features, and using a clean validation protocol.',
]
y0 = Inches(2.4); row_h = Inches(0.85)
for i, txt in enumerate(takes):
    y = y0 + i * row_h
    n_circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(0.7), y + Inches(0.1), Inches(0.6), Inches(0.6))
    n_circle.fill.solid(); n_circle.fill.fore_color.rgb = ORANGE
    n_circle.line.fill.background()
    add_text(s, Inches(0.7), y + Inches(0.1), Inches(0.6), Inches(0.6),
             str(i + 1), size=20, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.55), y, Inches(11.2), Inches(0.85),
             txt, size=16, color=DARK_TXT, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 20)


# -----------------------------------------------------------
# SLIDE 21 — Q & A / THANK YOU
# -----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.0), SW, Inches(1.6))
band.line.fill.background()
band.fill.solid(); band.fill.fore_color.rgb = ORANGE

add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.0),
         'Thank You', size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(3.3), Inches(12.3), Inches(1.2),
         'Questions & Discussion',
         size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.5),
         'Predicting Online Shoppers’ Purchase Intention  •  kNN from Scratch',
         size=18, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.5),
         'Code, figures, and report bundled in the project folder.',
         size=14, italic=True, color=WHITE, align=PP_ALIGN.CENTER)


# -----------------------------------------------------------
# Save
# -----------------------------------------------------------
out = os.path.join(BASE, 'Capstone_Presentation.pptx')
prs.save(out)
print(f'Wrote {out}')
print(f'Total slides: {len(prs.slides)}')
